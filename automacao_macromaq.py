import streamlit as st
import pandas as pd
from docx import Document
from docx.shared import Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
import io
import os
import zipfile
import unicodedata
import subprocess  # Para a conversão de PDF
import copy        # Para clonar as propriedades da linha se necessário
from datetime import datetime
from urllib.parse import quote
import base64

# --- CONFIGURAÇÕES ---
st.set_page_config(page_title="Automação SSMA Macromaq", layout="wide")

# Usar o diretório atual para facilitar o deploy no GitHub/Streamlit Cloud
BASE_PATH = os.getcwd()

# Caminhos originais na raiz do repositório
FUNDO_PATH = os.path.join(BASE_PATH, "fundo.png")
LOGO_PATH = os.path.join(BASE_PATH, "logo.png")
TEMPLATE_FICHA = os.path.join(BASE_PATH, "template_ficha.docx")
TEMPLATE_OS_JUNIOR = os.path.join(BASE_PATH, "template_os_Junior.docx")
TEMPLATE_NR06_JUNIOR = os.path.join(BASE_PATH, "template_nr06_Junior.pptx")
TEMPLATE_OS_SIMONE = os.path.join(BASE_PATH, "template_os_simone.docx")
TEMPLATE_NR06_SIMONE = os.path.join(BASE_PATH, "template_nr06_simone.pptx")

SHEET_ID = "1y98U3eK7JXJqQaMC0i7eFbwpvp97Nuyeml5Dis0UCUg"

# --- UNIDADES ---
UNIDADES = {
    "SÃO JOSÉ": {
        "CNPJ": "83.675.413/0001-01",
        "ENDERECO": "BR 101, km 210 / Bairro: Picadas do Sul – São José – SC / CEP: 88106-100"
    },
    "CHAPECÓ": {
        "CNPJ": "83.675.413/0002-84",
        "ENDERECO": "Rua Xanxerê, 360E – Bairro Líder – Chapecó/SC"
    },    
    "SÃO LEOPOLDO": {
        "CNPJ": "83.675.413/0016-80",
        "ENDERECO": "Avenida Senador Salgado Filho, 1970 – São Leopoldo – RS"
    },
    "JOINVILLE": {
        "CNPJ": "83.675.413/0011-75",
        "ENDERECO": "BR101, KM17 – Sentido Norte – Bairro Sta Catarina – Joinville / SC"
    },
    "PARANÁ": {
        "CNPJ": "83.675.413/0004-46",
        "ENDERECO": "Av. Juscelino K. de Oliveira, 3628 – Bairro CIC – Curitiba / PR"
    },
    "SÃO PAULO": {
        "CNPJ": "83.675.413/0008-70",
        "ENDERECO": "Rua Goiabeira 105/125 – Bairro Roseira de Cima – Jaguariúna / SP"
    },
    "MINAS GERAIS": {
        "CNPJ": "83.675.413/0014-18",
        "ENDERECO": "Anel Rodoviário Celso Mello Azevedo, 3713 - Bom Sucesso - BH/MG"
    },
    "ITAJAÍ": {
        "CNPJ": "83.675.413/0013-37",
        "ENDERECO": "Av. Vereador Abrahão João Francisco, 2300 - Dom Bosco - Itajaí / SC"
    }
}

# --- LAYOUT E CSS ---
def get_base64(bin_file):
    if not os.path.exists(bin_file):
        return ""
    with open(bin_file, 'rb') as f:
        data = f.read()
    return base64.b64encode(data).decode()

def aplicar_layout():
    try:
        fundo = get_base64(FUNDO_PATH)
        logo = get_base64(LOGO_PATH)
        st.markdown(f"""
        <style>
        /* REMOVER BARRA LATERAL E NAVEGAÇÃO NATIVA */
        [data-testid="stSidebar"], [data-testid="stSidebarNav"] {{
            display: none;
        }}

        .stApp {{
            background-image: url("data:image/png;base64,{fundo}");
            background-size: cover;
            background-position: center;
            background-attachment: fixed;
        }}
        .stSelectbox label, div[data-testid="stCheckbox"] label p {{
            color: white !important;
            background: rgba(0,0,0,0.7);
            padding: 5px 12px;
            border-radius: 8px;
            font-weight: bold;
        }}
        .stButton > button {{
            background: #2c3e50;
            color: #f9cc0b;
            border: 2px solid #f9cc0b;
            border-radius: 10px;
            height: 55px;
            font-weight: bold;
            width: 100%;
            font-size: 18px;
        }}
        .header-container {{
            display: flex;
            align-items: center;
            background: rgba(255,255,255,0.9);
            padding: 20px;
            border-radius: 15px;
            margin-bottom: 30px;
        }}
        .footer {{
            position: fixed;
            left: 0;
            bottom: 0;
            width: 100%;
            background: rgba(0,0,0,0.8);
            color: white;
            text-align: center;
            padding: 10px;
            font-size: 13px;
            z-index: 999;
        }}
        </style>
        <div class="header-container">
            <img src="data:image/png;base64,{logo}" width="320">
            <h1 style="margin-left:25px;color:#2c3e50;">Gestão SSMA Documentos</h1>
        </div>
        """, unsafe_allow_html=True)
    except:
        st.title("Automação SSMA Macromaq")

# --- FUNÇÕES AUXILIARES ---
def remover_acentos(texto):
    if not isinstance(texto, str): return str(texto)
    return "".join(c for c in unicodedata.normalize('NFD', texto.strip()) if unicodedata.category(c) != 'Mn').lower()

def limpar_valor(valor):
    if pd.isna(valor): 
        return ""
    return str(valor).strip()

def limpar_quebras_linha(texto):
    if not isinstance(texto, str): 
        texto = str(texto)
    texto = texto.replace('\r', ' ').replace('\n', ' ')
    while "  " in texto:
        texto = texto.replace("  ", " ")
    return texto.strip()

@st.cache_data(ttl=300)
def carregar_aba(aba_nome):
    try:
        url = f"https://docs.google.com/spreadsheets/d/{SHEET_ID}/gviz/tq?tqx=out:csv&sheet={quote(aba_nome.strip())}"
        df = pd.read_csv(url, dtype=str)
        return df
    except:
        return pd.DataFrame()

def data_extenso_pt():
    meses = {1:"Janeiro", 2:"Fevereiro", 3:"Março", 4:"Abril", 5:"Maio", 6:"Junho", 7:"Julho", 8:"Agosto", 9:"Setembro", 10:"Outubro", 11:"Novembro", 12:"Dezembro"}
    agora = datetime.now()
    return agora.strftime(f"%d de {meses[agora.month]} de %Y")

def formatar_matricula(valor):
    try: return str(int(float(valor))) if not pd.isna(valor) else ""
    except: return str(valor)

def formatar_cpf(cpf):
    try:
        if pd.isna(cpf):
            return "Não informado"
        cpf = str(cpf).strip()
        if cpf.endswith('.0'):
            cpf = cpf[:-2]
        if cpf.lower() in ["", "nan", "none", "0"]:
            return "Não informado"
        cpf = ''.join(filter(str.isdigit, cpf))
        if len(cpf) == 10:
            cpf = "0" + cpf
        if len(cpf) == 11:
            return f"{cpf[:3]}.{cpf[3:6]}.{cpf[6:9]}-{cpf[9:]}"
        return cpf if cpf != "" else "Não informado"
    except:
        return "Não informado"

def converter_para_pdf_linux(conteudo_arquivo, nome_original):
    try:
        temp_input = os.path.join(BASE_PATH, nome_original)
        with open(temp_input, "wb") as f:
            f.write(conteudo_arquivo)
        
        subprocess.run([
            'libreoffice', '--headless', '--convert-to', 'pdf', temp_input,
            '--outdir', BASE_PATH
        ], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
        
        nome_pdf = os.path.splitext(nome_original)[0] + ".pdf"
        temp_pdf_path = os.path.join(BASE_PATH, nome_pdf)
        
        if os.path.exists(temp_pdf_path):
            with open(temp_pdf_path, "rb") as f:
                pdf_bytes = f.read()
            
            os.remove(temp_input)
            os.remove(temp_pdf_path)
            return pdf_bytes, nome_pdf
    except Exception as e:
        st.sidebar.warning(f"Aviso técnico: Falha ao converter {nome_original} para PDF. Detalhes: {e}")
    return None, None

# --- PROCESSAMENTO DE DOCUMENTOS ---
def substituir_docx(doc, mapeamento):
    for p in doc.paragraphs:
        for tag, valor in mapeamento.items():
            if tag in p.text:
                for run in p.runs:
                    if tag in run.text:
                        run.text = run.text.replace(tag, str(valor))
                if tag in p.text:
                    texto_inteiro = p.text
                    for t, v in mapeamento.items():
                        texto_inteiro = texto_inteiro.replace(t, str(v))
                    if p.runs:
                        p.runs[0].text = texto_inteiro
                        for run in p.runs[1:]:
                            run.text = ""

    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    for tag, valor in mapeamento.items():
                        if tag in p.text:
                            for run in p.runs:
                                if tag in run.text:
                                    run.text = run.text.replace(tag, str(valor))
                            if tag in p.text:
                                texto_inteiro = p.text
                                for t, v in mapeamento.items():
                                    texto_inteiro = texto_inteiro.replace(t, str(v))
                                if p.runs:
                                    p.runs[0].text = texto_inteiro
                                    for run in p.runs[1:]:
                                        run.text = ""

def substituir_pptx(prs, mapeamento):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for tag, valor in mapeamento.items():
                            if tag in run.text: run.text = run.text.replace(tag, str(valor))

def preencher_ficha_docx(caminho_template, mapeamento, df_epis):
    doc = Document(caminho_template)
    substituir_docx(doc, mapeamento)
    
    tabela_alvo = None
    linhas_tags = []
    
    for tabela in doc.tables:
        for row in tabela.rows:
            texto_linha = "".join(cell.text for cell in row.cells)
            if "ITEM" in texto_linha or "DESC" in texto_linha:
                tabela_alvo = tabela
                linhas_tags.append(row)
                    
    if not tabela_alvo or len(linhas_tags) == 0:
        raise Exception("Nenhuma linha contendo a tag {{ITEM}} foi localizada no template do Word.")
        
    qtd_items = len(df_epis)
    linha_modelo = linhas_tags[0]
    
    def atualizar_celula_preservando_estilo(celula, novo_texto, alinhamento=WD_ALIGN_PARAGRAPH.CENTER):
        if not celula.paragraphs:
            celula.add_paragraph()
        p = celula.paragraphs[0]
        p.alignment = alinhamento
        if p.runs:
            p.runs[0].text = novo_texto
            for r in p.runs[1:]:
                r.text = ""
        else:
            p.add_run(novo_texto)

    for i, row_item in enumerate(linhas_tags):
        if i < qtd_items:
            item = df_epis.iloc[i]
            num_seq = f"{i + 1:02d}"
            
            for cell in row_item.cells:
                texto_celula = cell.text
                if "ITEM" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, num_seq, WD_ALIGN_PARAGRAPH.CENTER)
                elif "DESC" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('Descrição', '')), WD_ALIGN_PARAGRAPH.LEFT)
                elif "CA" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('C.A.', '')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "QT" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('qt.', '')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "unid" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('unid.', 'unid')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "DATA" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, datetime.now().strftime("%d/%m/%Y"), WD_ALIGN_PARAGRAPH.CENTER)
        else:
            for cell in row_item.cells:
                for paragraph in cell.paragraphs:
                    paragraph.text = ""

    if qtd_items > len(linhas_tags):
        tr_modelo = linha_modelo._tr
        for i in range(len(linhas_tags), qtd_items):
            item = df_epis.iloc[i]
            num_seq = f"{i + 1:02d}"
            
            nova_tr = copy.deepcopy(tr_modelo)
            nova_linha = tabela_alvo.add_row()
            nova_linha._tr.getparent().replace(nova_linha._tr, nova_tr)
            nova_linha._tr = nova_tr
            
            for cell in nova_linha.cells:
                texto_celula = cell.text
                if "ITEM" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, num_seq, WD_ALIGN_PARAGRAPH.CENTER)
                elif "DESC" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('Descrição', '')), WD_ALIGN_PARAGRAPH.LEFT)
                elif "CA" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('C.A.', '')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "QT" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('qt.', '')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "unid" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, limpar_valor(item.get('unid.', 'unid')), WD_ALIGN_PARAGRAPH.CENTER)
                elif "DATA" in texto_celula: 
                    atualizar_celula_preservando_estilo(cell, datetime.now().strftime("%d/%m/%Y"), WD_ALIGN_PARAGRAPH.CENTER)
            
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()

# --- APP LOGIC ---
aplicar_layout()

if st.button("🔄 Atualizar Dados da Planilha Agora"):
    st.cache_data.clear()
    st.rerun()

df_colab = carregar_aba("Colaboradores")
df_cargos = carregar_aba("Cargos")

if not df_colab.empty and not df_cargos.empty:
    df_colab['Nome_Formatado'] = df_colab['Nome Colaborador'].astype(str).str.strip().str.title()
    
    col1, col2, col3 = st.columns(3)
    with col1:
        lista_nomes = sorted(df_colab['Nome_Formatado'].dropna().unique())
        nome_sel = st.selectbox("1. Selecione o Colaborador:", lista_nomes)
        dados_colab = df_colab[df_colab['Nome_Formatado'] == nome_sel].iloc[0]
        
    with col2:
        unidade_plan = str(dados_colab.get('Filial', dados_colab.get('Unidade', ''))).upper().strip()
        lista_unid = list(UNIDADES.keys())
        idx = lista_unid.index(unidade_plan) if unidade_plan in lista_unid else 0
        unid_sel = st.selectbox("2. Unidade para OS:", lista_unid, index=idx)
    with col3:
        tecnico_sel = st.selectbox("3. Técnico Responsável:", ["Técnico Junior", "Técnica Simone"])

    t_os = TEMPLATE_OS_JUNIOR if tecnico_sel == "Técnico Junior" else TEMPLATE_OS_SIMONE
    t_nr = TEMPLATE_NR06_JUNIOR if tecnico_sel == "Técnico Junior" else TEMPLATE_NR06_SIMONE

    st.markdown("<br>", unsafe_allow_html=True)
    c1, c2, c3 = st.columns(3)
    g_os, g_ficha, g_cert = c1.checkbox("OS", True), c2.checkbox("Ficha EPI", True), c3.checkbox("Certificado", True)

    incluir_pdf = st.checkbox("📄 Incluir cópias em formato PDF no Kit", True)

    if st.button("🚀 PROCESSAR DOCUMENTOS"):
        with st.spinner("Gerando documentos..."):
            cargo = str(dados_colab['Cargo']).strip()
            arquivos = {}
            df_cargos['f_l'] = df_cargos['Função'].astype(str).apply(remover_acentos)
            desc_f = df_cargos[df_cargos['f_l'] == remover_acentos(cargo)]

            if not desc_f.empty:
                desc_atv = limpar_quebras_linha(desc_f['Descrição da Atividade'].fillna('').values[0]) if 'Descrição da Atividade' in desc_f.columns else ''
                riscos_agentes = limpar_quebras_linha(desc_f['Riscos e Agentes Existentes'].fillna('').values[0]) if 'Riscos e Agentes Existentes' in desc_f.columns else ''
                medidas_protecao = limpar_quebras_linha(desc_f['Medidas de Proteção'].fillna('').values[0]) if 'Medidas de Proteção' in desc_f.columns else ''
                
                setor_original = limpar_valor(dados_colab.get('NomeLocal', dados_colab.get('Setor', '')))
                setor_final = setor_original if setor_original != "" else unid_sel.title()
                
                coluna_cpf = [col for col in df_colab.columns if 'CPF' in col.upper()]
                if coluna_cpf:
                    cpf_bruto = dados_colab[coluna_cpf[0]]
                else:
                    try:
                        cpf_bruto = dados_colab.iloc[18]
                    except:
                        cpf_bruto = ""
                        
                cpf_final = formatar_cpf(cpf_bruto)
                
                # 1. Ordem de Serviço
                if g_os:
                    doc = Document(t_os)
                    substituir_docx(doc, {
                        "{{NOME}}": dados_colab['Nome Colaborador'], 
                        "{{FUNCAO}}": cargo.upper(), 
                        "{{CNPJ}}": UNIDADES[unid_sel]["CNPJ"], 
                        "{{ENDERECO}}": UNIDADES[unid_sel]["ENDERECO"], 
                        "{{SETOR}}": setor_final, 
                        "{{DESCRICAO_ATIVIDADE}}": desc_atv, 
                        "{{MEDIDAS_PROTECAO}}": medidas_protecao,    
                        "{{RISCOS_AGENTES}}": riscos_agentes,
                        "{{DATA}}": datetime.now().strftime("%d/%m/%Y")
                    })
                    b = io.BytesIO(); doc.save(b)
                    conteudo_docx = b.getvalue()
                    nome_docx = f"OS {nome_sel}.docx"
                    arquivos[nome_docx] = conteudo_docx
                    
                    if incluir_pdf:
                        pdf_bytes, nome_pdf = converter_para_pdf_linux(conteudo_docx, nome_docx)
                        if pdf_bytes: arquivos[nome_pdf] = pdf_bytes

                # 2. Ficha de EPI
                if g_ficha:
                    cargo_limpo = cargo.strip()
                    df_e = carregar_aba(cargo_limpo)
                    
                    if df_e.empty:
                        df_e = carregar_aba(f"{cargo_limpo} ")
                    if df_e.empty:
                        df_e = carregar_aba(cargo_limpo.title())
                    if df_e.empty: 
                        df_e = carregar_aba(remover_acentos(cargo_limpo))
                    if df_e.empty and "jr" in cargo_limpo.lower():
                        df_e = carregar_aba(cargo_limpo.lower().replace("jr", "Jr"))

                    if not df_e.empty:
                        m_f = {"{{NOME}}": dados_colab['Nome Colaborador'], "{{MATRICULA}}": formatar_matricula(dados_colab.get('Matrícula', '')), "{{FUNCAO}}": cargo, "{{DATA_ADMISSAO}}": datetime.now().strftime("%d/%m/%Y"), "{{SETOR}}": setor_final, "{{CENTRO_CUSTO}}": ""}
                        conteudo_ficha_docx = preencher_ficha_docx(TEMPLATE_FICHA, m_f, df_e)
                        nome_ficha_docx = f"Ficha EPI {nome_sel}.docx"
                        arquivos[nome_ficha_docx] = conteudo_ficha_docx
                        
                        if incluir_pdf:
                            pdf_bytes, nome_pdf = converter_para_pdf_linux(conteudo_ficha_docx, nome_ficha_docx)
                            if pdf_bytes: arquivos[nome_pdf] = pdf_bytes
                    else:
                        st.error(f"❌ Erro crítico: A aba de EPIs para o cargo '{cargo_limpo}' não pôde ser baixada.")

                # 3. Certificado NR06
                if g_cert:
                    prs = Presentation(t_nr)
                    
                    if "SÃO JOSÉ" in unid_sel.upper():
                        local_data_string = f"{data_extenso_pt()}."
                    else:
                        local_data_string = f"{unid_sel.title()}, {data_extenso_pt()}."

                    substituir_pptx(prs, {
                        "{{NOME}}": dados_colab['Nome Colaborador'], 
                        "{{CPF}}": cpf_final, 
                        "{{FUNCAO}}": cargo, 
                        "{{DATA_TREINAMENTO}}": datetime.now().strftime("%d/%m/%Y"), 
                        "{{LOCAL_DATA}}": local_data_string
                    })
                    b = io.BytesIO(); prs.save(b)
                    conteudo_pptx = b.getvalue()
                    nome_pptx = f"NR06 {nome_sel}.pptx"
                    arquivos[nome_pptx] = conteudo_pptx
                    
                    if incluir_pdf:
                        pdf_bytes, nome_pdf = converter_para_pdf_linux(conteudo_pptx, nome_pptx)
                        if pdf_bytes: arquivos[nome_pdf] = pdf_bytes

                if arquivos:
                    z_b = io.BytesIO()
                    with zipfile.ZipFile(z_b, "w") as z:
                        for n, d in arquivos.items(): z.writestr(n, d)
                    st.success("✅ Documentos prontos!")
                    st.download_button("📦 BAIXAR KIT COMPLETO (ZIP)", z_b.getvalue(), f"Kit_{nome_sel}.zip", use_container_width=True)
            else:
                st.error(f"Cargo '{cargo}' não encontrado na aba Cargos.")
else:
    st.error("Erro ao carregar dados da planilha Google. Verifique o acesso público ou as abas.")

# --- FOOTER ---
st.markdown("""<div class="footer">© 2026 Gestão Documentos | Desenvolvido por: Dilceu Junior</div>""", unsafe_allow_html=True)
